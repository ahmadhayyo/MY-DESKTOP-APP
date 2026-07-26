"""
powerpoint_tools.py — Full PowerPoint (.pptx) creation & editing.

Built on python-pptx (no Microsoft Office required — works headless, never hangs).
Covers a real office workflow: build a deck, add title / bullet / image / table /
chart slides, read & extract content, edit text, restyle, and export to PDF.

Tools:
  • ppt_create(path, title, subtitle)                 — new deck with a title slide
  • ppt_add_slide(path, layout, title, content)       — add a slide (many layouts)
  • ppt_add_bullets(path, title, bullets)             — title + bulleted content slide
  • ppt_add_image(path, image_path, title, caption)   — image slide
  • ppt_add_table(path, title, data)                  — table slide from JSON
  • ppt_add_chart(path, title, chart_type, data)      — bar/line/pie chart slide
  • ppt_read(path)                                     — extract all text per slide
  • ppt_edit_text(path, find, replace)                — find/replace across the deck
  • ppt_set_theme(path, primary_hex, font)            — recolor titles + set font
  • ppt_to_pdf(path, pdf_path)                        — export to PDF (uses Office if present)
  • ppt_info(path)                                    — slide count, titles, dimensions
"""
from __future__ import annotations

import json
import os
from pathlib import Path

from langchain_core.tools import tool


# ── helpers ───────────────────────────────────────────────────────────────────
def _require_pptx():
    try:
        import pptx  # noqa: F401
        return None
    except ImportError:
        return ("❌ مكتبة python-pptx غير مثبّتة.\n"
                "   ثبّتها بـ: pip install python-pptx")


def _hex_to_rgb(hex_str: str):
    from pptx.dml.color import RGBColor
    h = (hex_str or "").lstrip("#").strip()
    if len(h) != 6:
        h = "1F4E79"  # default office blue
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _open_or_err(path: str):
    """Return (Presentation, None) or (None, error_str)."""
    from pptx import Presentation
    p = Path(path)
    if not p.is_file():
        return None, f"❌ الملف غير موجود: {path}\n   أنشئه أولاً بـ ppt_create()."
    try:
        return Presentation(str(p)), None
    except Exception as e:
        return None, f"❌ تعذّر فتح العرض: {e}"


# ═══════════════════════════════════════════════════════════════════════════════
#  CREATE
# ═══════════════════════════════════════════════════════════════════════════════
@tool
def ppt_create(
    path: str,
    title: str = "",
    subtitle: str = "",
) -> str:
    """إنشاء عرض PowerPoint جديد (.pptx) مع شريحة عنوان.

    Args:
        path: مسار الملف، مثل C:/Users/PT/Desktop/عرض.pptx
        title: عنوان الشريحة الأولى
        subtitle: نص فرعي (اختياري)

    بعد الإنشاء، أضف شرائح بـ ppt_add_bullets / ppt_add_table / ppt_add_chart / ppt_add_image.
    """
    err = _require_pptx()
    if err:
        return err
    from pptx import Presentation
    from pptx.util import Inches
    try:
        prs = Presentation()
        # Force modern 16:9 widescreen (python-pptx default is legacy 4:3).
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
        if slide.shapes.title is not None:
            slide.shapes.title.text = title or "عرض تقديمي"
        # subtitle placeholder (idx 1)
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                ph.text = subtitle
                break
        p = Path(path)
        p.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(p))
        return f"✅ تم إنشاء العرض: {p}\n   شريحة العنوان: «{title}»\n   أضف المزيد بـ ppt_add_bullets()/ppt_add_table()/ppt_add_chart()."
    except Exception as e:
        return f"❌ خطأ في إنشاء العرض: {e}"


# ═══════════════════════════════════════════════════════════════════════════════
#  ADD SLIDES
# ═══════════════════════════════════════════════════════════════════════════════
@tool
def ppt_add_slide(
    path: str,
    title: str = "",
    content: str = "",
    layout: str = "title_content",
) -> str:
    """إضافة شريحة عامة إلى عرض موجود.

    Args:
        path: مسار العرض الموجود
        title: عنوان الشريحة
        content: نص المحتوى (أسطر متعددة = نقاط منفصلة)
        layout: 'title_content' (عنوان+محتوى) أو 'title_only' أو 'blank' أو 'section'
    """
    prs, err = _open_or_err(path)
    if err:
        return err
    try:
        layout_map = {"title_content": 1, "section": 2, "blank": 6, "title_only": 5}
        idx = layout_map.get(layout, 1)
        idx = min(idx, len(prs.slide_layouts) - 1)
        slide = prs.slides.add_slide(prs.slide_layouts[idx])
        if title and slide.shapes.title is not None:
            slide.shapes.title.text = title
        if content:
            body = None
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    body = ph
                    break
            if body is not None:
                tf = body.text_frame
                lines = [l for l in content.split("\n")]
                tf.text = lines[0] if lines else ""
                for line in lines[1:]:
                    para = tf.add_paragraph()
                    para.text = line
        prs.save(path)
        return f"✅ أُضيفت شريحة «{title}» (التخطيط: {layout}). إجمالي الشرائح: {len(prs.slides)}"
    except Exception as e:
        return f"❌ خطأ في إضافة الشريحة: {e}"


@tool
def ppt_add_bullets(
    path: str,
    title: str,
    bullets: str,
) -> str:
    """إضافة شريحة بعنوان وقائمة نقاط (الأكثر شيوعاً في العروض).

    Args:
        path: مسار العرض
        title: عنوان الشريحة
        bullets: النقاط — سطر لكل نقطة. ابدأ السطر بمسافتين للنقاط الفرعية (تداخل).
                 مثال:
                   الإيرادات ارتفعت 20%
                   التكاليف انخفضت
                     في قسم التسويق
                     في قسم اللوجستيات
    """
    prs, err = _open_or_err(path)
    if err:
        return err
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide.shapes.title is not None:
            slide.shapes.title.text = title
        body = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                body = ph
                break
        if body is None:
            return "❌ تخطيط الشريحة لا يحتوي مكان محتوى."
        tf = body.text_frame
        tf.clear()
        lines = [l.rstrip() for l in bullets.split("\n") if l.strip()]
        first = True
        for line in lines:
            # leading spaces → nesting level
            stripped = line.lstrip(" ")
            level = min((len(line) - len(stripped)) // 2, 4)
            if first:
                tf.text = stripped
                tf.paragraphs[0].level = level
                first = False
            else:
                para = tf.add_paragraph()
                para.text = stripped
                para.level = level
        prs.save(path)
        return f"✅ أُضيفت شريحة نقاط «{title}» ({len(lines)} نقطة)."
    except Exception as e:
        return f"❌ خطأ: {e}"


@tool
def ppt_add_image(
    path: str,
    image_path: str,
    title: str = "",
    caption: str = "",
) -> str:
    """إضافة شريحة تحتوي صورة (مع عنوان وتعليق اختياريين).

    Args:
        path: مسار العرض
        image_path: مسار الصورة (png/jpg)
        title: عنوان أعلى الشريحة
        caption: تعليق أسفل الصورة
    """
    prs, err = _open_or_err(path)
    if err:
        return err
    if not Path(image_path).is_file():
        return f"❌ الصورة غير موجودة: {image_path}"
    try:
        from pptx.util import Inches, Pt
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
        if title and slide.shapes.title is not None:
            slide.shapes.title.text = title
        # Fit image within the slide, centered
        sw, sh = prs.slide_width, prs.slide_height
        top = Inches(1.5)
        max_w = sw - Inches(2)
        pic = slide.shapes.add_picture(image_path, Inches(1), top, width=max_w)
        # center horizontally
        pic.left = int((sw - pic.width) / 2)
        if caption:
            box = slide.shapes.add_textbox(Inches(1), sh - Inches(1), sw - Inches(2), Inches(0.6))
            box.text_frame.text = caption
            box.text_frame.paragraphs[0].font.size = Pt(14)
        prs.save(path)
        return f"✅ أُضيفت شريحة صورة «{title or image_path}»."
    except Exception as e:
        return f"❌ خطأ في إضافة الصورة: {e}"


@tool
def ppt_add_table(
    path: str,
    title: str,
    data: str,
) -> str:
    """إضافة شريحة تحتوي جدولاً من بيانات JSON.

    Args:
        path: مسار العرض
        title: عنوان الشريحة
        data: JSON — قائمة قوائم (أول صف = العناوين) أو قائمة قواميس.
              مثال: [["المنتج","الكمية","السعر"],["قلم",100,2],["دفتر",50,5]]
    """
    prs, err = _open_or_err(path)
    if err:
        return err
    try:
        rows_data = json.loads(data) if isinstance(data, str) else data
        if not rows_data:
            return "❌ لا توجد بيانات للجدول."
        # normalize dicts → rows
        if isinstance(rows_data[0], dict):
            headers = list(rows_data[0].keys())
            table_rows = [headers] + [[r.get(h, "") for h in headers] for r in rows_data]
        else:
            table_rows = [list(r) for r in rows_data]

        from pptx.util import Inches, Pt
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        if title and slide.shapes.title is not None:
            slide.shapes.title.text = title
        n_rows, n_cols = len(table_rows), len(table_rows[0])
        gfx = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(0.5), Inches(1.6),
            prs.slide_width - Inches(1), Inches(0.4 * n_rows),
        )
        table = gfx.table
        for r, row in enumerate(table_rows):
            for c in range(n_cols):
                val = row[c] if c < len(row) else ""
                cell = table.cell(r, c)
                cell.text = str(val)
                for para in cell.text_frame.paragraphs:
                    para.font.size = Pt(14)
                    if r == 0:
                        para.font.bold = True
        prs.save(path)
        return f"✅ أُضيفت شريحة جدول «{title}» ({n_rows}×{n_cols})."
    except json.JSONDecodeError:
        return "❌ بيانات JSON غير صحيحة."
    except Exception as e:
        return f"❌ خطأ في الجدول: {e}"


@tool
def ppt_add_chart(
    path: str,
    title: str,
    chart_type: str,
    data: str,
) -> str:
    """إضافة شريحة تحتوي رسماً بيانياً.

    Args:
        path: مسار العرض
        title: عنوان الشريحة
        chart_type: 'bar' (أعمدة) أو 'line' (خطي) أو 'pie' (دائري)
        data: JSON بالشكل:
              {"categories": ["يناير","فبراير","مارس"],
               "series": {"المبيعات": [100,120,90], "الأرباح": [30,40,25]}}
              للرسم الدائري استخدم سلسلة واحدة فقط.
    """
    prs, err = _open_or_err(path)
    if err:
        return err
    try:
        spec = json.loads(data) if isinstance(data, str) else data
        categories = spec.get("categories", [])
        series = spec.get("series", {})
        if not categories or not series:
            return "❌ يجب توفير categories و series."

        from pptx.util import Inches
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        type_map = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
        }
        xl_type = type_map.get(chart_type.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        if title and slide.shapes.title is not None:
            slide.shapes.title.text = title

        chart_data = CategoryChartData()
        chart_data.categories = categories
        for name, values in series.items():
            chart_data.add_series(name, values)

        x, y = Inches(1), Inches(1.6)
        cx, cy = prs.slide_width - Inches(2), prs.slide_height - Inches(2.5)
        slide.shapes.add_chart(xl_type, x, y, cx, cy, chart_data)
        prs.save(path)
        return f"✅ أُضيفت شريحة رسم بياني «{title}» (نوع: {chart_type}, {len(series)} سلسلة)."
    except json.JSONDecodeError:
        return "❌ بيانات JSON غير صحيحة."
    except Exception as e:
        return f"❌ خطأ في الرسم البياني: {e}"


# ═══════════════════════════════════════════════════════════════════════════════
#  READ / EDIT / STYLE
# ═══════════════════════════════════════════════════════════════════════════════
@tool
def ppt_read(path: str) -> str:
    """قراءة كل النصوص من عرض PowerPoint، شريحة بشريحة."""
    prs, err = _open_or_err(path)
    if err:
        return err
    try:
        out = [f"📊 العرض: {Path(path).name} ({len(prs.slides)} شريحة)\n"]
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs) or para.text
                        if t.strip():
                            texts.append(t.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        texts.append(" | ".join(c.text for c in row.cells))
            out.append(f"── شريحة {i} ──")
            out.append("\n".join(texts) if texts else "(لا نص)")
            out.append("")
        return "\n".join(out)
    except Exception as e:
        return f"❌ خطأ في القراءة: {e}"


@tool
def ppt_edit_text(path: str, find: str, replace: str) -> str:
    """البحث عن نص واستبداله عبر كل شرائح العرض (يحافظ على التنسيق قدر الإمكان)."""
    prs, err = _open_or_err(path)
    if err:
        return err
    try:
        count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if find in run.text:
                            run.text = run.text.replace(find, replace)
                            count += 1
        prs.save(path)
        return (f"✅ تم استبدال «{find}» بـ «{replace}» في {count} موضع."
                if count else f"ℹ️ لم يُعثر على «{find}».")
    except Exception as e:
        return f"❌ خطأ في التعديل: {e}"


@tool
def ppt_set_theme(path: str, primary_hex: str = "1F4E79", font: str = "") -> str:
    """تطبيق لون أساسي على كل عناوين الشرائح، وخط موحّد اختياري.

    Args:
        path: مسار العرض
        primary_hex: لون العناوين بصيغة hex مثل '1F4E79' أو 'C00000'
        font: اسم الخط (اختياري) مثل 'Calibri' أو 'Arial'
    """
    prs, err = _open_or_err(path)
    if err:
        return err
    try:
        color = _hex_to_rgb(primary_hex)
        changed = 0
        for slide in prs.slides:
            title = slide.shapes.title
            if title is not None and title.has_text_frame:
                for para in title.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = color
                        run.font.bold = True
                        if font:
                            run.font.name = font
                    if not para.runs and para.text:
                        para.font.color.rgb = color
                        para.font.bold = True
                changed += 1
            if font:
                for shape in slide.shapes:
                    if shape.has_text_frame and shape != title:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                run.font.name = font
        prs.save(path)
        return f"✅ طُبّق اللون #{primary_hex} على {changed} عنوان" + (f" والخط «{font}»" if font else "") + "."
    except Exception as e:
        return f"❌ خطأ في التنسيق: {e}"


@tool
def ppt_info(path: str) -> str:
    """معلومات سريعة: عدد الشرائح، عناوينها، أبعاد العرض."""
    prs, err = _open_or_err(path)
    if err:
        return err
    try:
        from pptx.util import Emu
        slides = list(prs.slides)
        lines = [f"📊 {Path(path).name}",
                 f"   الشرائح: {len(slides)}",
                 f"   الأبعاد: {Emu(prs.slide_width).inches:.1f}×{Emu(prs.slide_height).inches:.1f} بوصة"]
        for i, s in enumerate(slides, 1):
            t = s.shapes.title.text if s.shapes.title is not None and s.shapes.title.has_text_frame else "(بلا عنوان)"
            lines.append(f"   {i}. {t}")
        return "\n".join(lines)
    except Exception as e:
        return f"❌ خطأ: {e}"


@tool
def ppt_to_pdf(path: str, pdf_path: str = "") -> str:
    """تحويل عرض PowerPoint إلى PDF.

    يستخدم Microsoft PowerPoint عبر COM إن كان مثبّتاً (أعلى جودة)،
    وإلا يبلّغ بأن التحويل يتطلب PowerPoint.
    """
    src = Path(path)
    if not src.is_file():
        return f"❌ الملف غير موجود: {path}"
    out = Path(pdf_path) if pdf_path else src.with_suffix(".pdf")
    out.parent.mkdir(parents=True, exist_ok=True)
    try:
        import win32com.client  # type: ignore
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        try:
            deck = powerpoint.Presentations.Open(str(src.resolve()), WithWindow=False)
            deck.SaveAs(str(out.resolve()), 32)  # 32 = ppSaveAsPDF
            deck.Close()
            return f"✅ تم التصدير إلى PDF: {out}"
        finally:
            powerpoint.Quit()
    except Exception as e:
        return (f"⚠️ تعذّر التحويل عبر PowerPoint COM ({e}).\n"
                f"   تأكد أن Microsoft PowerPoint مثبّت، أو افتح العرض واحفظه PDF يدوياً.")
