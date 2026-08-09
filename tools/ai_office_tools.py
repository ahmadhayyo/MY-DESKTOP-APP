"""
AI Office Tools — write what you want, send your data, get the exact file.

For humanitarian-relief reporting (World Vision, North-Syria camps): the user
types a plain instruction and pastes/points at data; a strong model (Claude
Sonnet by default, free fallback) structures it, and these tools render the
precise deliverable — Excel, Word, or PowerPoint — deterministically.

  • ai_to_excel        — instruction + data → a clean, formatted spreadsheet
  • ai_to_word         — instruction + data → a professional Word report
  • ai_to_powerpoint   — instruction + data → a presentation deck
  • ai_office          — auto-picks the format from `output`

The model supplies intelligence; the rendering is done by code, so numbers and
Arabic land exactly. If `data` is a file path (.csv/.xlsx/.txt/.docx/.pdf) its
content is read first, so she can point at a source file instead of pasting.
"""
from __future__ import annotations

import os
from pathlib import Path
from typing import Annotated

from langchain_core.tools import tool

from config import resolve_output_path
from core import ai_office as _ai
from core import office_convert as _oc


def _load_data(data: str) -> str:
    """If `data` is a path to a supported file, load a text/tabular view of it;
    otherwise treat `data` as the raw pasted content."""
    s = (data or "").strip()
    if len(s) < 500 and os.path.exists(os.path.expanduser(s)):
        p = os.path.expanduser(s)
        ext = Path(p).suffix.lower()
        try:
            if ext in (".xlsx", ".xlsm"):
                rows = _oc.xlsx_to_rows(p)
            elif ext in (".csv", ".tsv"):
                rows = _oc.csv_to_rows(p)
            elif ext == ".txt":
                return Path(p).read_text(encoding="utf-8", errors="replace")
            elif ext == ".docx":
                tabs = _oc.docx_tables_to_rows(p)
                if tabs:
                    rows = tabs[0]
                else:
                    from docx import Document
                    return "\n".join(par.text for par in Document(p).paragraphs)
            elif ext == ".pdf":
                tabs = _oc.pdf_tables_to_rows(p)
                rows = tabs[0] if tabs else []
                if not rows:
                    from pypdf import PdfReader
                    return "\n".join((pg.extract_text() or "") for pg in PdfReader(p).pages)
            else:
                return data
            return "\n".join("\t".join(str(c) for c in r) for r in rows)
        except Exception:
            return data
    return data


# ── Excel ─────────────────────────────────────────────────────────────────────
@tool
def ai_to_excel(
    instruction: Annotated[str, "ماذا تريد أن ينفّذ النموذج بالبيانات (بالعربية أو "
                                "الإنجليزية): تنظيم/تصنيف/حساب/فرز/ترجمة أعمدة... إلخ."],
    data: Annotated[str, "البيانات الخام (منسوخة) أو مسار ملف مصدر."],
    dest: Annotated[str, "مسار ملف الإكسل الناتج (.xlsx)."],
) -> str:
    """AI-structure raw data into a precise, formatted Excel per your instruction.
    A strong model organises/cleans/computes; the file is rendered exactly."""
    try:
        content = _load_data(data)
        spec = _ai.plan_table(instruction, content)
        header, rows = spec.get("header", []), spec.get("rows", [])
        all_rows = ([header] + rows) if header else rows
        if not all_rows:
            return "❌ لم يُنتج النموذج أي صفوف."
        out = resolve_output_path(dest)
        path = _oc.rows_to_xlsx(all_rows, out, sheet_name=spec.get("sheet_name") or "Data",
                                header=bool(header), rtl=spec.get("rtl"),
                                title=spec.get("title", ""))
        return (f"✅ تم إنشاء الإكسل بدقّة: {path}\n"
                f"📊 {len(rows)} صف × {len(header) or (len(rows[0]) if rows else 0)} عمود · "
                f"النموذج: {spec.get('_model', '?')}")
    except Exception as exc:
        return f"❌ ai_to_excel: {exc}"


# ── Word ──────────────────────────────────────────────────────────────────────
def _render_docx(spec: dict, out: str) -> str:
    from docx import Document
    from docx.shared import Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    rtl = bool(spec.get("rtl"))
    doc = Document()

    def _rtl_align(par):
        if rtl:
            par.alignment = WD_ALIGN_PARAGRAPH.RIGHT

    if spec.get("title"):
        _rtl_align(doc.add_heading(spec["title"], level=0))
    for sec in spec.get("sections", []):
        if sec.get("heading"):
            _rtl_align(doc.add_heading(sec["heading"], level=1))
        for para in sec.get("paragraphs", []) or []:
            _rtl_align(doc.add_paragraph(str(para)))
        for b in sec.get("bullets", []) or []:
            _rtl_align(doc.add_paragraph(str(b), style="List Bullet"))
        tbl = sec.get("table")
        if tbl and (tbl.get("header") or tbl.get("rows")):
            rows = ([tbl["header"]] if tbl.get("header") else []) + (tbl.get("rows") or [])
            if rows:
                cols = max(len(r) for r in rows)
                t = doc.add_table(rows=0, cols=cols)
                t.style = "Light Grid Accent 1"
                if rtl:
                    from docx.oxml.ns import qn
                    t._tbl.tblPr.append(t._tbl.tblPr.makeelement(qn("w:bidiVisual"), {}))
                for i, r in enumerate(rows):
                    cells = t.add_row().cells
                    for j in range(cols):
                        cells[j].text = str(r[j]) if j < len(r) else ""
                        for p in cells[j].paragraphs:
                            if rtl:
                                p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                            for run in p.runs:
                                run.font.size = Pt(11)
                                if tbl.get("header") and i == 0:
                                    run.font.bold = True
    Path(out).parent.mkdir(parents=True, exist_ok=True)
    doc.save(out)
    return out


@tool
def ai_to_word(
    instruction: Annotated[str, "ما المطلوب كتابته/تنظيمه في تقرير Word."],
    data: Annotated[str, "البيانات الخام أو مسار ملف مصدر."],
    dest: Annotated[str, "مسار ملف Word الناتج (.docx)."],
) -> str:
    """AI-build a professional Word report (headings, paragraphs, bullets, tables)
    from your instruction and data — RTL-aware for Arabic."""
    try:
        content = _load_data(data)
        spec = _ai.plan_document(instruction, content)
        out = resolve_output_path(dest)
        path = _render_docx(spec, out)
        return (f"✅ تم إنشاء تقرير Word: {path}\n"
                f"🧩 {len(spec.get('sections', []))} قسم · النموذج: {spec.get('_model', '?')}")
    except Exception as exc:
        return f"❌ ai_to_word: {exc}"


# ── PowerPoint ────────────────────────────────────────────────────────────────
def _render_pptx(spec: dict, out: str) -> str:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation()
    rtl = bool(spec.get("rtl"))

    # Title slide
    title_layout = prs.slide_layouts[0]
    s0 = prs.slides.add_slide(title_layout)
    s0.shapes.title.text = spec.get("title") or "تقرير"
    for sec in spec.get("sections", []):
        layout = prs.slide_layouts[1]  # title + content
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = sec.get("heading") or ""
        body = slide.placeholders[1].text_frame
        body.clear()
        first = True
        items = list(sec.get("paragraphs", []) or []) + list(sec.get("bullets", []) or [])
        for it in items:
            p = body.paragraphs[0] if first else body.add_paragraph()
            p.text = str(it)
            p.font.size = Pt(18)
            first = False
        tbl = sec.get("table")
        if tbl and (tbl.get("header") or tbl.get("rows")):
            rows = ([tbl["header"]] if tbl.get("header") else []) + (tbl.get("rows") or [])
            if rows:
                nrows, ncols = len(rows), max(len(r) for r in rows)
                gtbl = slide.shapes.add_table(nrows, ncols, Inches(0.5), Inches(2.2),
                                              Inches(9), Inches(0.4 * nrows)).table
                for i, r in enumerate(rows):
                    for j in range(ncols):
                        gtbl.cell(i, j).text = str(r[j]) if j < len(r) else ""
    Path(out).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    return out


@tool
def ai_to_powerpoint(
    instruction: Annotated[str, "ما المطلوب في العرض التقديمي."],
    data: Annotated[str, "البيانات الخام أو مسار ملف مصدر."],
    dest: Annotated[str, "مسار ملف PowerPoint الناتج (.pptx)."],
) -> str:
    """AI-build a PowerPoint deck (title slide + a slide per section with bullets
    and tables) from your instruction and data."""
    try:
        content = _load_data(data)
        spec = _ai.plan_document(instruction, content)
        out = resolve_output_path(dest)
        path = _render_pptx(spec, out)
        return (f"✅ تم إنشاء العرض: {path}\n"
                f"🖼️ {len(spec.get('sections', [])) + 1} شريحة · النموذج: {spec.get('_model', '?')}")
    except Exception as exc:
        return f"❌ ai_to_powerpoint: {exc}"


# ── Auto-dispatch ─────────────────────────────────────────────────────────────
@tool
def ai_office(
    instruction: Annotated[str, "ماذا تريد أن ينفّذ النموذج."],
    data: Annotated[str, "البيانات الخام أو مسار ملف مصدر."],
    output: Annotated[str, "نوع الملف الناتج: excel / word / powerpoint."],
    dest: Annotated[str, "مسار الناتج (مع الامتداد المناسب)."],
) -> str:
    """One entry point: describe the task, send data, choose output — Excel, Word,
    or PowerPoint — and get the exact file, AI-structured and code-rendered."""
    o = (output or "").strip().lower()
    if o in ("excel", "xlsx", "اكسل", "إكسل", "اكسيل"):
        return ai_to_excel.invoke({"instruction": instruction, "data": data, "dest": dest})
    if o in ("word", "docx", "وورد", "ورد"):
        return ai_to_word.invoke({"instruction": instruction, "data": data, "dest": dest})
    if o in ("powerpoint", "pptx", "ppt", "بوربوينت", "باوربوينت", "عرض"):
        return ai_to_powerpoint.invoke({"instruction": instruction, "data": data, "dest": dest})
    return f"❌ نوع الملف غير مدعوم: {output} — استخدم excel / word / powerpoint."
